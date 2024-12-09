from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from loguru import logger
import os

class PPTXCreator:
    def __init__(self):
        self.prs = Presentation()
        #self.prs = Presentation("templates/MasterTemplate.pptx")
        self.output_dir = "outputs"
        if not os.path.exists(self.output_dir):
            os.makedirs(self.output_dir)

    def add_title_slide(self, title, subtitle=""):
        """添加标题页"""
        title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = subtitle

    def add_content_slide(self, title, content):
        """添加内容页"""
        bullet_slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(bullet_slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title
        
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        
        for item in content:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0

    def add_image_slide(self, title, image_path):
        """添加图片页"""
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)
        
        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_box.text = title
        
        # 添加图片
        if os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(2), Inches(2), width=Inches(6))
        else:
            logger.error(f"图片文件不存在: {image_path}")

    def add_chart_slide(self, title, categories, values):
        """添加图表页"""
        chart_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(chart_slide_layout)
        
        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_box.text = title
        
        # 创建图表数据
        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series('数据系列1', values)
        
        # 添加图表
        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        )

    def save(self, filename="presentation.pptx"):
        """保存PPT文件"""
        output_path = os.path.join(self.output_dir, filename)
        self.prs.save(output_path)
        logger.info(f"PPT文件已保存: {output_path}")

def main():
    # 创建示例PPT
    creator = PPTXCreator()
    
    # 添加标题页
    creator.add_title_slide("演示文稿标题", "副标题")
    
    # 添加内容页
    content = ["第一点", "第二点", "第三点"]
    creator.add_content_slide("内容页面", content)
    
    # 添加图片页
    creator.add_image_slide("图片展示", "images/performance_chart.png")
    
    # 添加图表页
    categories = ['A', 'B', 'C', 'D']
    values = [4.3, 2.5, 3.5, 4.5]
    creator.add_chart_slide("数据图表", categories, values)
    
    # 保存文件
    creator.save()

if __name__ == "__main__":
    main()
